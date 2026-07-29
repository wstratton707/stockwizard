"""
pptx_builder.py — Professional PowerPoint report builder for QuantWizard
Generates a polished stock analysis or portfolio deck using python-pptx + matplotlib.
"""

import io
import os
import math
import re
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
from datetime import datetime
from constants import get_risk_free_rate

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    # Fallbacks so module-level constants below don't raise NameError when
    # python-pptx isn't installed. The builder functions check PPTX_AVAILABLE
    # before doing real work, so these stand-ins are never used for output.
    def RGBColor(*_a, **_k): return None
    def Inches(v):           return v
    def Pt(v):               return v
    def Emu(v):              return v
    PP_ALIGN = None


# ── Brand colours ─────────────────────────────────────────────────────────────
C_NAVY      = RGBColor(0x1F, 0x4E, 0x79)   # dark navy — headers, accents
C_BLUE      = RGBColor(0x2E, 0x75, 0xB6)   # mid blue — sub-headers
C_ACCENT    = RGBColor(0x00, 0xB0, 0xF0)   # bright cyan — highlights
C_GREEN     = RGBColor(0x70, 0xAD, 0x47)   # positive values
C_RED       = RGBColor(0xC0, 0x00, 0x00)   # negative values
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xF0, 0xF7, 0xFF)   # light blue bg
C_DARK_TEXT = RGBColor(0x1E, 0x29, 0x3B)
C_GREY_TEXT = RGBColor(0x64, 0x74, 0x8B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# Full logo (light line-art on transparent) — shows on the dark cover slide.
_ASSET_LOGO = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets", "logo_full.png")

MPL_COLORS = {
    "navy":   "#1F4E79",
    "blue":   "#2E75B6",
    "cyan":   "#00B0F0",
    "green":  "#70AD47",
    "red":    "#C00000",
    "orange": "#E8A838",
    "purple": "#8E44AD",
    "grey":   "#94A3B8",
}


# ── Slide helpers ─────────────────────────────────────────────────────────────

def _new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # completely blank
    return prs.slides.add_slide(blank_layout)


def _rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def _text_box(slide, text, l, t, w, h,
              font_size=12, bold=False, italic=False,
              color=C_DARK_TEXT, align=PP_ALIGN.LEFT,
              font_name="Calibri"):   # cleaner, modern Office default (was Arial)
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size    = Pt(font_size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    run.font.name    = font_name
    return txb


def _add_image(slide, buf, l, t, w, h):
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(l), Inches(t), Inches(w), Inches(h))


def _slide_header(slide, title, subtitle=None):
    """Dark navy top bar with title."""
    _rect(slide, 0, 0, 13.33, 1.15, fill_rgb=C_NAVY)
    _text_box(slide, title, 0.35, 0.12, 11.5, 0.65,
              font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        _text_box(slide, subtitle, 0.35, 0.72, 11.5, 0.35,
                  font_size=11, italic=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    # thin accent line below header
    _rect(slide, 0, 1.15, 13.33, 0.04, fill_rgb=C_ACCENT)


def _slide_footer(slide, page_num, total_pages):
    _rect(slide, 0, 7.2, 13.33, 0.3, fill_rgb=C_NAVY)
    _text_box(slide, "QuantWizard  |  For informational purposes only. Not investment advice.",
              0.3, 7.22, 10, 0.25, font_size=8, color=C_GREY_TEXT)
    _text_box(slide, f"{page_num} / {total_pages}",
              12.5, 7.22, 0.7, 0.25, font_size=8, color=C_GREY_TEXT, align=PP_ALIGN.RIGHT)


def _kv_block(slide, pairs, l, t, w, col_w=1.8, row_h=0.52, bg=True):
    """Render a list of (label, value, positive?) tuples as a metric table."""
    for i, item in enumerate(pairs):
        label, value = item[0], item[1]
        positive     = item[2] if len(item) > 2 else None
        row_t = t + i * row_h
        if bg:
            bg_col = C_LIGHT if i % 2 == 0 else C_WHITE
            _rect(slide, l, row_t, w, row_h - 0.02, fill_rgb=bg_col)
        _text_box(slide, label, l + 0.1, row_t + 0.07, col_w, row_h - 0.1,
                  font_size=10, color=C_GREY_TEXT)
        val_color = C_DARK_TEXT
        if positive is True:
            val_color = C_GREEN
        elif positive is False:
            val_color = C_RED
        _text_box(slide, str(value), l + col_w, row_t + 0.07, w - col_w - 0.1, row_h - 0.1,
                  font_size=11, bold=True, color=val_color, align=PP_ALIGN.RIGHT)


# ── Matplotlib chart helpers ──────────────────────────────────────────────────

def _chart_style(ax, title, xlabel="Date", ylabel=""):
    ax.set_title(title, fontsize=12, fontweight="bold", color="#1F4E79", pad=8)
    ax.set_xlabel(xlabel, fontsize=9, color="#475569")
    ax.set_ylabel(ylabel, fontsize=9, color="#475569")
    ax.tick_params(axis="x", rotation=30, labelsize=7.5)
    ax.tick_params(axis="y", labelsize=8)
    ax.grid(axis="y", linestyle="--", alpha=0.35, color="#CBD5E1")
    ax.grid(axis="x", linestyle=":",  alpha=0.2,  color="#CBD5E1")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.spines["left"].set_color("#CBD5E1")
    ax.spines["bottom"].set_color("#CBD5E1")
    ax.set_facecolor("#F8FAFC")
    ax.figure.patch.set_facecolor("#FFFFFF")
    if ax.get_legend_handles_labels()[0]:
        ax.legend(fontsize=8.5, framealpha=0.8, edgecolor="#E2E8F0")


def _fig_to_buf(fig, dpi=140):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=dpi, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    buf.seek(0)
    return buf


def _price_ma_chart(df, ticker, w=11, h=4.2):
    dates = pd.to_datetime(df["Date"])
    fig, ax = plt.subplots(figsize=(w, h))
    ax.plot(dates, df["Close"], color=MPL_COLORS["navy"], linewidth=1.8,
            label="Close", zorder=3)
    for ma, col, lw in [("MA20", MPL_COLORS["orange"], 1.1),
                         ("MA50", MPL_COLORS["green"],  1.1),
                         ("MA200", MPL_COLORS["red"],   1.1)]:
        if ma in df.columns and df[ma].notna().any():
            ax.plot(dates, df[ma], color=col, linewidth=lw,
                    linestyle="--", label=ma, zorder=2)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"${x:,.2f}"))
    _chart_style(ax, f"{ticker} — Price & Moving Averages", ylabel="Price ($)")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _volume_chart(df, ticker, w=11, h=2.8):
    dates  = pd.to_datetime(df["Date"])
    fig, ax = plt.subplots(figsize=(w, h))
    if "Volume" in df.columns:
        daily_ret = df.get("Daily_Return", pd.Series([0]*len(df))).fillna(0)
        colors = [MPL_COLORS["green"] if r >= 0 else MPL_COLORS["red"] for r in daily_ret]
        ax.bar(dates, df["Volume"], color=colors, width=1.2, alpha=0.7)
        if "Vol_MA20" in df.columns:
            ax.plot(dates, df["Vol_MA20"], color=MPL_COLORS["navy"],
                    linewidth=1.3, linestyle="--", label="20-Day Avg")
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x/1e6:.1f}M"))
    elif "Vol_MA20" in df.columns:
        ax.plot(dates, df["Vol_MA20"], color=MPL_COLORS["navy"],
                linewidth=1.4, label="20-Day Avg Volume")
        ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x/1e6:.1f}M"))
    _chart_style(ax, f"{ticker} — Volume", ylabel="Volume")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _bollinger_chart(df, ticker, w=11, h=4.2):
    if "BB_Upper" not in df.columns:
        return None
    dates = pd.to_datetime(df["Date"])
    fig, ax = plt.subplots(figsize=(w, h))
    ax.plot(dates, df["Close"],     color=MPL_COLORS["navy"],  linewidth=1.8, label="Close",    zorder=3)
    ax.plot(dates, df["BB_Upper"],  color=MPL_COLORS["red"],   linewidth=1.0, linestyle="--", label="BB Upper")
    ax.plot(dates, df["BB_Middle"], color=MPL_COLORS["grey"],  linewidth=1.0, linestyle="--", label="BB Middle")
    ax.plot(dates, df["BB_Lower"],  color=MPL_COLORS["green"], linewidth=1.0, linestyle="--", label="BB Lower")
    ax.fill_between(dates, df["BB_Upper"], df["BB_Lower"],
                    alpha=0.07, color=MPL_COLORS["blue"])
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"${x:,.2f}"))
    _chart_style(ax, f"{ticker} — Bollinger Bands (20-day, 2σ)", ylabel="Price ($)")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _rsi_chart(df, ticker, w=11, h=2.8):
    if "RSI14" not in df.columns:
        return None
    dates = pd.to_datetime(df["Date"])
    fig, ax = plt.subplots(figsize=(w, h))
    ax.plot(dates, df["RSI14"], color=MPL_COLORS["purple"], linewidth=1.4, label="RSI (14)")
    ax.axhline(70, color=MPL_COLORS["red"],   linewidth=0.8, linestyle="--", alpha=0.7, label="Overbought (70)")
    ax.axhline(30, color=MPL_COLORS["green"], linewidth=0.8, linestyle="--", alpha=0.7, label="Oversold (30)")
    ax.fill_between(dates, df["RSI14"], 70,
                    where=df["RSI14"] >= 70, alpha=0.15, color=MPL_COLORS["red"])
    ax.fill_between(dates, df["RSI14"], 30,
                    where=df["RSI14"] <= 30, alpha=0.15, color=MPL_COLORS["green"])
    ax.set_ylim(0, 100)
    _chart_style(ax, f"{ticker} — RSI (14)", ylabel="RSI")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _cumulative_chart(df, ticker, w=11, h=4.2):
    if "Cumulative_Index" not in df.columns:
        return None
    dates = pd.to_datetime(df["Date"])
    fig, ax = plt.subplots(figsize=(w, h))
    ax.plot(dates, df["Cumulative_Index"], color=MPL_COLORS["navy"],
            linewidth=2.0, label=ticker, zorder=3)
    bench_colors = [MPL_COLORS["red"], MPL_COLORS["green"],
                    MPL_COLORS["orange"], MPL_COLORS["purple"]]
    cum_cols = [c for c in df.columns if c.endswith("_Cumulative")]
    for i, col in enumerate(cum_cols):
        label = col.replace("_Cumulative", "")
        ax.plot(dates, df[col], color=bench_colors[i % len(bench_colors)],
                linewidth=1.3, linestyle="--", label=label)
    ax.axhline(100, color="#CBD5E1", linewidth=0.7, linestyle=":")
    ax.fill_between(dates, df["Cumulative_Index"], 100,
                    where=df["Cumulative_Index"] >= 100,
                    alpha=0.08, color=MPL_COLORS["green"])
    ax.fill_between(dates, df["Cumulative_Index"], 100,
                    where=df["Cumulative_Index"] < 100,
                    alpha=0.08, color=MPL_COLORS["red"])
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:.0f}"))
    _chart_style(ax, f"{ticker} — Cumulative Return vs Benchmarks",
                 ylabel="Index (100 = start)")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _monte_carlo_chart(mc_sim_df, mc_summary, ticker, w=11, h=4.2):
    if mc_sim_df is None:
        return None
    pct_labels = ["P5 (Bear)", "P25 (Low)", "P50 (Median)", "P75 (Bull)", "P95 (Best)"]
    pct_colors = [MPL_COLORS["red"], MPL_COLORS["orange"], MPL_COLORS["navy"],
                  MPL_COLORS["green"], "#27AE60"]
    # Plot the FULL forecast horizon (was capped at 252, which truncated the
    # portfolio's multi-year forecast to just year one). For long horizons show
    # the x-axis in years; for ≤1yr (e.g. the single-stock deck) keep days.
    n     = len(mc_sim_df)
    arr   = mc_sim_df.values[:n]
    pcts  = np.percentile(arr, [5, 25, 50, 75, 95], axis=1)
    years = n / 252.0
    if years > 1.5:
        x_vals, xlabel, horizon = [d / 252.0 for d in range(n)], "Years Forward", f"{round(years)}-Year"
    else:
        x_vals, xlabel, horizon = list(range(n)), "Trading Days Forward", f"{n} Trading Days"
    fig, ax = plt.subplots(figsize=(w, h))
    ax.fill_between(x_vals, pcts[0], pcts[4], alpha=0.10, color=MPL_COLORS["blue"], label="P5–P95 range")
    ax.fill_between(x_vals, pcts[1], pcts[3], alpha=0.18, color=MPL_COLORS["blue"], label="P25–P75 range")
    for j, (lbl, col) in enumerate(zip(pct_labels, pct_colors)):
        lw = 2.2 if "Median" in lbl else 1.0
        ls = "-"  if "Median" in lbl else "--"
        ax.plot(x_vals, pcts[j], color=col, linewidth=lw, linestyle=ls, label=lbl)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"${x:,.0f}"))
    _chart_style(ax, f"{ticker} — Monte Carlo Forecast ({horizon})",
                 xlabel=xlabel, ylabel="Value ($)")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _drawdown_chart(df, ticker, w=11, h=2.8):
    if "Drawdown_60d" not in df.columns:
        return None
    dates = pd.to_datetime(df["Date"])
    fig, ax = plt.subplots(figsize=(w, h))
    dd = df["Drawdown_60d"] * 100
    ax.fill_between(dates, dd, 0, alpha=0.5, color=MPL_COLORS["red"], label="Drawdown (%)")
    ax.plot(dates, dd, color=MPL_COLORS["red"], linewidth=0.8)
    ax.axhline(-20, color="#888", linewidth=0.7, linestyle="--", alpha=0.5, label="-20% threshold")
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:.1f}%"))
    _chart_style(ax, f"{ticker} — 60-Day Rolling Drawdown", ylabel="Drawdown (%)")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _valuation_range_chart(dcf, ticker, w=11.6, h=1.9):
    """Football field: the bear→bull fair-value band with the base case and
    today's price marked on it. One picture of the whole valuation argument."""
    scn = (dcf or {}).get("scenarios") or {}
    fvs = {k: (scn.get(k) or {}).get("fair_value") for k in ("bear", "base", "bull")}
    if not all(isinstance(v, (int, float)) for v in fvs.values()):
        return None
    price = dcf.get("price")
    if not isinstance(price, (int, float)) or price <= 0:
        return None

    lo, hi = min(fvs["bear"], fvs["bull"]), max(fvs["bear"], fvs["bull"])
    fig, ax = plt.subplots(figsize=(w, h))
    ax.barh([0], hi - lo, left=lo, height=0.34, color="#D6E4F0",
            edgecolor=MPL_COLORS["blue"], linewidth=1.1, zorder=2)
    ax.scatter([fvs["base"]], [0], marker="D", s=130, color=MPL_COLORS["navy"],
               zorder=5, label="DCF base case")
    ax.axvline(price, color=MPL_COLORS["red"], linewidth=2.0, linestyle="--",
               zorder=4, label="Today's price")

    # Band ends read outwards so they can never collide with each other or with
    # the base-case label, however narrow the bear→bull range is.
    ax.annotate(f"Bear ${lo:,.0f}", (lo, 0), xytext=(-8, 0), textcoords="offset points",
                ha="right", va="center", fontsize=8.5, color="#475569")
    ax.annotate(f"Bull ${hi:,.0f}", (hi, 0), xytext=(8, 0), textcoords="offset points",
                ha="left", va="center", fontsize=8.5, color="#475569")
    ax.annotate(f"Base ${fvs['base']:,.0f}", (fvs["base"], 0), xytext=(0, 16),
                textcoords="offset points", ha="center", fontsize=9,
                fontweight="bold", color=MPL_COLORS["navy"])
    # White plaque so the dashed price line doesn't run through its own label.
    ax.annotate(f"Price ${price:,.0f}", (price, 0), xytext=(0, 34),
                textcoords="offset points", ha="center", fontsize=9,
                fontweight="bold", color=MPL_COLORS["red"],
                bbox=dict(boxstyle="round,pad=0.28", facecolor="#FFFFFF",
                          edgecolor=MPL_COLORS["red"], linewidth=0.8))

    span = max(hi - lo, 1e-9)
    pad = max(span * 0.45, abs(price - fvs["base"]) * 0.30, price * 0.08)
    ax.set_xlim(min(lo, price) - pad, max(hi, price) + pad)
    ax.set_ylim(-0.55, 0.80)
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"${x:,.0f}"))
    _chart_style(ax, f"{ticker} — Fair Value Range vs Today's Price",
                 xlabel="Value per share", ylabel="")
    ax.tick_params(axis="x", rotation=0, labelsize=8.5)
    ax.set_yticks([])
    ax.grid(axis="y", visible=False)
    ax.grid(axis="x", linestyle="--", alpha=0.35, color="#CBD5E1")
    ax.spines["left"].set_visible(False)
    fig.tight_layout()
    return _fig_to_buf(fig)


# ── Portfolio chart helpers ───────────────────────────────────────────────────

def _alloc_pie_chart(weights, ticker_info, w=6, h=4.5):
    # Sort once and use the SAME ordering for wedges and the legend — previously
    # the legend zipped sorted sizes against the dict's unsorted keys, so every
    # ticker was labelled with the wrong slice's percentage.
    ordered = sorted(weights.items(), key=lambda x: -x[1])   # [(ticker, weight), ...] desc
    sizes   = [wt * 100 for _, wt in ordered]
    colors  = plt.cm.Blues(np.linspace(0.85, 0.40, len(sizes)))   # largest = darkest
    fig, ax = plt.subplots(figsize=(w, h))
    wedges, texts, autotexts = ax.pie(
        sizes, labels=None, colors=colors,
        autopct=lambda p: f"{p:.1f}%" if p > 3 else "",
        startangle=140, pctdistance=0.78,
        wedgeprops={"linewidth": 1.2, "edgecolor": "white"}
    )
    for at in autotexts:
        at.set_fontsize(8)
        at.set_color("white")
        at.set_fontweight("bold")
    ax.legend(wedges, [f"{tk} ({wt * 100:.1f}%)" for tk, wt in ordered],
              loc="center left", bbox_to_anchor=(1, 0, 0.5, 1),
              fontsize=8, framealpha=0.7)
    ax.set_title("Portfolio Allocation", fontsize=12, fontweight="bold",
                 color="#1F4E79", pad=6)
    fig.patch.set_facecolor("#FFFFFF")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _portfolio_performance_chart(backtest_df, w=11, h=4.0):
    """Account value over the backtest: Portfolio vs a same-schedule S&P 500
    benchmark, with the cumulative contributions line for context. backtest_df is
    indexed by date with columns Portfolio / Contrib / NAV / SP500."""
    if backtest_df is None or backtest_df.empty:
        return None
    dates = pd.to_datetime(backtest_df.index)
    series = [
        ("Portfolio", "Portfolio",          MPL_COLORS["navy"],   2.0, "-"),
        ("SP500",     "S&P 500 (same DCA)",  MPL_COLORS["grey"],   1.5, "--"),
        ("Contrib",   "Total Contributed",   MPL_COLORS["orange"], 1.3, ":"),
    ]
    fig, ax = plt.subplots(figsize=(w, h))
    plotted = False
    for col, label, clr, lw, ls in series:
        if col in backtest_df.columns and not backtest_df[col].isna().all():
            ax.plot(dates, backtest_df[col], color=clr, linewidth=lw, linestyle=ls, label=label)
            plotted = True
    if not plotted:
        plt.close(fig)
        return None
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"${x:,.0f}"))
    _chart_style(ax, "Portfolio vs Benchmark — Account Value",
                 ylabel="Account Value ($)")
    fig.tight_layout()
    return _fig_to_buf(fig)


def _holdings_bar_chart(stock_metrics, w=9, h=4.0):
    if not stock_metrics:
        return None
    tickers = list(stock_metrics.keys())
    # ann_return is already stored as a percentage (compute_stock_metrics), so no ×100.
    returns = [stock_metrics[t].get("ann_return", 0) for t in tickers]
    colors  = [MPL_COLORS["green"] if r >= 0 else MPL_COLORS["red"] for r in returns]
    fig, ax = plt.subplots(figsize=(w, h))
    bars = ax.bar(tickers, returns, color=colors, alpha=0.85, edgecolor="white", linewidth=1.2)
    for bar, val in zip(bars, returns):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + (0.3 if val >= 0 else -1.2),
                f"{val:+.1f}%", ha="center", va="bottom", fontsize=8.5, fontweight="bold",
                color="#1E293B")
    ax.axhline(0, color="#94A3B8", linewidth=0.8)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f"{x:.0f}%"))
    _chart_style(ax, "Annualized Return by Holding", xlabel="", ylabel="Ann. Return (%)")
    fig.tight_layout()
    return _fig_to_buf(fig)


# ── Stock Deck ────────────────────────────────────────────────────────────────

# Plain-English translation of analysis.dcf_valuation()'s failure reasons, so a
# missing model reads as an explanation rather than a blank slide.
_DCF_NO_MODEL_REASONS = {
    "fundamentals unavailable":
        "No company financial statements are published for this security — normal for "
        "ETFs, funds, ADRs and crypto — so there is no cash-flow stream to discount.",
    "no price":
        "No current market price was available, so the model has nothing to anchor to.",
    "no market cap":
        "No market capitalisation was available, so the share count could not be derived.",
    "WACC must exceed terminal growth":
        "The discount rate is not above the assumed terminal growth rate, which makes the "
        "terminal value mathematically infinite.",
    "no positive free cash flow to project":
        "This company has not reported positive free cash flow in the years available, so "
        "there is nothing to project forward. Judge it on multiples and quality instead.",
}


def _build_valuation_slide(prs, ticker, dcf, page_num, total):
    """The signature analytic, as a conclusion slide: what today's price already
    assumes, what the model says it's worth, and the range around that."""
    s = _blank_slide(prs)

    if not isinstance(dcf, dict) or not dcf.get("ok"):
        _slide_header(s, "Valuation", ticker)
        _slide_footer(s, page_num, total)
        reason = dcf.get("reason", "") if isinstance(dcf, dict) else ""
        _text_box(s, "No discounted-cash-flow model for this security.",
                  0.6, 2.7, 12.1, 0.5, font_size=17, bold=True,
                  color=C_NAVY, align=PP_ALIGN.CENTER)
        _text_box(s, _DCF_NO_MODEL_REASONS.get(
                      reason,
                      "The inputs a DCF needs — audited financial statements, a market "
                      "price and a positive free-cash-flow history — were not all "
                      "available for this security."),
                  1.9, 3.35, 9.5, 1.1, font_size=12.5,
                  color=C_GREY_TEXT, align=PP_ALIGN.CENTER)
        return True

    imp   = dcf.get("market_implied_growth")
    price = dcf["price"]
    up    = dcf.get("upside")
    scn   = dcf.get("scenarios") or {}

    _slide_header(s, "Valuation — What Today's Price Implies",
                  f"{ticker}  ·  Reverse discounted cash flow  ·  "
                  f"{dcf['wacc']*100:.1f}% discount rate  ·  {dcf['years']}-year horizon")
    _slide_footer(s, page_num, total)

    # ── Hero: the implied-growth statement ───────────────────────────────────
    _rect(s, 0.3, 1.38, 12.7, 1.66, fill_rgb=C_LIGHT)
    _rect(s, 0.3, 1.38, 0.09, 1.66, fill_rgb=C_ACCENT)
    if imp is not None:
        _text_box(s, f"{imp*100:.1f}%", 0.6, 1.62, 3.0, 0.95,
                  font_size=48, bold=True, color=C_NAVY)
        _text_box(s, "IMPLIED FCF GROWTH, PER YEAR", 0.63, 2.55, 3.2, 0.3,
                  font_size=9, bold=True, color=C_GREY_TEXT)
        _hero = (f"To justify ${price:,.2f} a share, {ticker}'s free cash flow would have "
                 f"to compound at about {imp*100:.1f}% a year for {dcf['years']} years, "
                 f"then {dcf['terminal_growth']*100:.1f}% in perpetuity.\n"
                 f"This model's own base case assumes {dcf['base_growth']*100:.1f}%. "
                 f"The question is not what we think the stock is worth — it is whether "
                 f"you believe this company can clear that bar.")
    else:
        _text_box(s, "n/a", 0.6, 1.62, 3.0, 0.95,
                  font_size=48, bold=True, color=C_GREY_TEXT)
        _text_box(s, "IMPLIED FCF GROWTH, PER YEAR", 0.63, 2.55, 3.2, 0.3,
                  font_size=9, bold=True, color=C_GREY_TEXT)
        _hero = (f"Today's ${price:,.2f} price sits outside the range of growth rates this "
                 f"model can solve for, so no single implied rate can be quoted.\n"
                 f"That is itself a finding: at a {dcf['wacc']*100:.1f}% discount rate the "
                 f"price is not explicable by free-cash-flow growth alone.")
    _text_box(s, _hero, 4.05, 1.60, 8.75, 1.30, font_size=12.5, color=C_DARK_TEXT)

    # ── Conclusion tiles ─────────────────────────────────────────────────────
    _tiles = [
        ("Base-case fair value", f"${dcf['fair_value']:,.2f}" if dcf.get("fair_value") else "N/A", None),
        ("Today's price",        f"${price:,.2f}", None),
        ("Upside / downside",    f"{up*100:+.1f}%" if up is not None else "N/A",
                                 (up > 0) if up is not None else None),
        ("Verdict",              ("Undervalued" if up is not None and up > 0.15 else
                                  "Overvalued"  if up is not None and up < -0.15 else
                                  "Fairly valued"),
                                 (up > 0.15) if up is not None and abs(up) > 0.15 else None),
    ]
    _tw = 3.07
    for _i, (_lbl, _val, _pos) in enumerate(_tiles):
        _x = 0.3 + _i * (_tw + 0.18)
        _rect(s, _x, 3.22, _tw, 1.32, fill_rgb=C_WHITE,
              line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width=1)
        _rect(s, _x, 3.22, _tw, 0.07, fill_rgb=C_ACCENT)
        _vcol = C_GREEN if _pos is True else (C_RED if _pos is False else C_NAVY)
        _text_box(s, _lbl.upper(), _x + 0.18, 3.42, _tw - 0.3, 0.28,
                  font_size=9, bold=True, color=C_GREY_TEXT)
        _text_box(s, _val, _x + 0.18, 3.76, _tw - 0.3, 0.62,
                  font_size=(22 if len(str(_val)) > 9 else 25), bold=True, color=_vcol)

    # ── Football field ───────────────────────────────────────────────────────
    _chart = None
    try:
        _chart = _valuation_range_chart(dcf, ticker)
    except Exception:
        _chart = None
    if _chart is not None:
        _add_image(s, _chart, 0.32, 4.66, 12.66, 2.07)
    else:
        _rows = [(n.title(),
                  f"${(scn.get(n) or {}).get('fair_value', 0):,.2f}"
                  if (scn.get(n) or {}).get("fair_value") else "N/A",
                  ((scn.get(n) or {}).get("upside") or 0) > 0)
                 for n in ("bear", "base", "bull")]
        _kv_block(s, _rows, 0.32, 4.75, 6.2, col_w=2.6, row_h=0.5)

    _text_box(s, f"Two-stage unlevered DCF on free cash flow · {dcf['wacc']*100:.1f}% discount "
                 f"rate · {dcf['terminal_growth']*100:.1f}% terminal growth · base FCF "
                 f"normalised over recent years · full editable model, assumptions and "
                 f"sensitivity grid on the Valuation sheet of the Excel workbook. "
                 f"A model, not a price target.",
              0.35, 6.80, 12.6, 0.36, font_size=8.5, italic=True, color=C_GREY_TEXT)
    return True


def _build_fundamentals_slide(prs, ticker, fundamentals, page_num, total):
    """One slide of EDGAR-sourced fundamentals + quality scores. Returns False if
    no fundamentals are available (e.g. crypto)."""
    f = fundamentals
    if not f or not f.get("ok"):
        return False
    v, m, r, l, g = f["valuation"], f["margins"], f["returns"], f["leverage"], f["growth"]
    q, fc = f.get("quality", {}), f.get("fcf", {})

    def pc(x, suffix=""):
        return f"{x}{suffix}" if x is not None else "N/A"

    s = _blank_slide(prs)
    _slide_header(s, f"{ticker} — Fundamentals & Valuation",
                  f"Source: {f.get('source','—')}   ·   FY ending {f.get('as_of','—')}")

    left = [
        ("P/E", pc(v["pe"], "x")), ("P/S", pc(v["ps"], "x")), ("P/B", pc(v["pb"], "x")),
        ("EV / EBITDA", pc(f.get("ev_ebitda"), "x")), ("FCF Yield", pc(fc.get("fcf_yield"), "%")),
        ("Gross Margin", pc(m["gross"], "%")), ("Operating Margin", pc(m["operating"], "%")),
        ("Net Margin", pc(m["net"], "%")), ("Return on Equity", pc(r["roe"], "%")),
    ]
    _fs, _z, _zone = q.get("f_score"), q.get("z_score"), q.get("z_zone")
    right = [
        ("Revenue YoY", pc(g["revenue_yoy"], "%"), (g["revenue_yoy"] or 0) >= 0),
        ("EPS YoY", pc(g["eps_yoy"], "%"), (g["eps_yoy"] or 0) >= 0),
        ("Revenue CAGR", pc(g["revenue_cagr"], "%")), ("EPS CAGR", pc(g["eps_cagr"], "%")),
        ("Current Ratio", pc(l["current_ratio"])), ("Debt / Equity", pc(l["debt_to_equity"])),
        ("Free Cash Flow", f"${fc['fcf']/1e9:,.1f}B" if fc.get("fcf") is not None else "N/A"),
        ("Piotroski F-Score", f"{_fs} / 9" if _fs is not None else "N/A",
         (_fs >= 7) if _fs is not None else None),
        ("Altman Z-Score", f"{_z} ({_zone})" if _z is not None else "N/A",
         (_zone == "safe") if _z is not None else None),
    ]
    _text_box(s, "Valuation & Profitability", 0.5, 1.4, 5.9, 0.35,
              font_size=13, bold=True, color=C_NAVY)
    _kv_block(s, left, 0.5, 1.85, 5.95, col_w=3.1, row_h=0.52)
    _text_box(s, "Growth & Quality", 6.9, 1.4, 5.9, 0.35,
              font_size=13, bold=True, color=C_NAVY)
    _kv_block(s, right, 6.9, 1.85, 5.95, col_w=3.1, row_h=0.52)
    _slide_footer(s, page_num, total)
    return True


def build_stock_pptx(ticker, df, period_label,
                     company_details=None, mc_sim_df=None, mc_summary=None,
                     news_list=None, summary_text="", fundamentals=None,
                     dcf=None):
    """Build a professional stock analysis PowerPoint. Returns BytesIO."""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed.")

    prs        = _new_prs()
    cd         = company_details or {}
    latest     = df.iloc[-1]
    first      = df.iloc[0]
    period_ret = (latest["Close"] / first["Close"] - 1) * 100
    ret        = df["Daily_Return"].dropna()
    ann_ret    = ret.mean() * 252
    ann_std    = ret.std() * np.sqrt(252)
    downside   = ret[ret < 0].std() * np.sqrt(252)
    # Excess-return Sharpe/Sortino (subtract the risk-free rate) so the exported
    # deck matches the on-screen metric cards and the portfolio engine.
    rfr        = get_risk_free_rate()
    sharpe     = (ann_ret - rfr) / ann_std  if ann_std  else float("nan")
    sortino    = (ann_ret - rfr) / downside if downside else float("nan")
    max_dd     = df["Drawdown_60d"].min() * 100 if "Drawdown_60d" in df.columns else float("nan")

    def _fmt_pct(v):
        if v is None or (isinstance(v, float) and math.isnan(v)):
            return "N/A"
        return f"{v:+.2f}%"

    def _fmt_ratio(v):
        if v is None or (isinstance(v, float) and math.isnan(v)):
            return "N/A"
        return f"{v:.2f}"

    # +1 for the "Bottom Line" executive-summary slide added after the cover, and
    # +1 for the valuation slide whenever a DCF was attempted at all (it renders a
    # short explanation rather than nothing when the model couldn't be built).
    total_slides = ((11 if (fundamentals and fundamentals.get("ok")) else 10)
                    + 1 + (1 if dcf is not None else 0))
    page = [1]   # cover is page 1 (hardcoded); next slide starts at 2

    def _next_page():
        page[0] += 1
        return page[0]

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill_rgb=C_NAVY)
    _rect(sl, 0, 5.5, 13.33, 2.0, fill_rgb=RGBColor(0x0B, 0x14, 0x24))   # KPI band
    _rect(sl, 0, 5.5, 13.33, 0.03, fill_rgb=C_ACCENT)                    # accent divider
    _rect(sl, 0.6, 1.35, 0.09, 2.45, fill_rgb=C_ACCENT)                  # left accent bar

    company_name = cd.get("Name", ticker)
    # Full logo top-right on the dark cover (replaces the plain "QUANTWIZARD"
    # wordmark — the logo carries the name). Falls back to the wordmark text.
    if os.path.exists(_ASSET_LOGO):
        try:
            sl.shapes.add_picture(_ASSET_LOGO, Inches(10.05), Inches(0.55), height=Inches(2.05))
        except Exception:
            _text_box(sl, "QUANTWIZARD", 0.9, 0.85, 12, 0.5,
                      font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    else:
        _text_box(sl, "QUANTWIZARD", 0.9, 0.85, 12, 0.5,
                  font_size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)
    _text_box(sl, ticker, 0.86, 1.3, 12, 1.2,
              font_size=66, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    _text_box(sl, company_name, 0.9, 2.65, 12, 0.6,
              font_size=22, bold=False, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.LEFT)
    _text_box(sl, "Equity Research Report", 0.9, 3.35, 12, 0.45,
              font_size=15, italic=True, color=C_ACCENT, align=PP_ALIGN.LEFT)

    sector   = cd.get("Sector", "")
    exchange = cd.get("Exchange", "")
    meta_str = "   ·   ".join(filter(None, [sector.title() if sector else "", exchange, f"Period: {period_label}"]))
    _text_box(sl, meta_str, 0.9, 3.9, 12, 0.35,
              font_size=11, color=RGBColor(0x94, 0xA3, 0xB8), align=PP_ALIGN.LEFT)

    _text_box(sl, f"Generated {datetime.now().strftime('%B %d, %Y')}   ·   "
                  f"Multi-source data: Polygon · Yahoo Finance · Finnhub · SEC EDGAR",
              0.9, 4.45, 12, 0.3, font_size=9.5, color=C_GREY_TEXT)
    _text_box(sl, "For informational purposes only. Not financial advice.",
              0.9, 4.78, 12, 0.3, font_size=8.5, italic=True, color=RGBColor(0x64, 0x74, 0x8B))

    # KPI stat band — mirrors the web hero ribbon for an immediate read
    _kpis = [
        ("CURRENT PRICE",   f"${latest['Close']:,.2f}",       C_WHITE),
        ("PERIOD RETURN",   f"{period_ret:+.1f}%",            C_GREEN if period_ret >= 0 else C_RED),
        ("SHARPE RATIO",    _fmt_ratio(sharpe),               C_WHITE),
        ("ANN. VOLATILITY", f"{ann_std * 100:.1f}%",          C_WHITE),
    ]
    for _i, (_lbl, _val, _col) in enumerate(_kpis):
        _x = 0.9 + _i * 3.05
        _text_box(sl, _val, _x, 5.85, 2.9, 0.75, font_size=30, bold=True, color=_col, align=PP_ALIGN.LEFT)
        _text_box(sl, _lbl, _x + 0.03, 6.65, 2.9, 0.3, font_size=9.5, bold=True, color=C_ACCENT, align=PP_ALIGN.LEFT)

    _text_box(sl, "1", 12.9, 7.15, 0.4, 0.25, font_size=8, color=C_GREY_TEXT, align=PP_ALIGN.RIGHT)

    # ── Slide 2: The Bottom Line (executive summary) ──────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "The Bottom Line", ticker)
    _slide_footer(sl, _next_page(), total_slides)

    _rect(sl, 0.3, 1.45, 12.7, 2.15, fill_rgb=C_LIGHT)
    _rect(sl, 0.3, 1.45, 0.09, 2.15, fill_rgb=C_ACCENT)
    _text_box(sl, "VERDICT", 0.55, 1.62, 11, 0.3, font_size=11, bold=True, color=C_NAVY)
    _bl_text = summary_text or f"{ticker} analysis generated by QuantWizard."
    if len(_bl_text) > 620:
        _bl_text = _bl_text[:620].rsplit(" ", 1)[0] + "…"
    _text_box(sl, _bl_text, 0.55, 1.98, 12.2, 1.55, font_size=12, color=C_DARK_TEXT)

    _bl_kpis = [
        ("Current Price", f"${latest['Close']:,.2f}", None),
        ("Period Return", f"{period_ret:+.1f}%",      period_ret >= 0),
        ("Sharpe Ratio",  _fmt_ratio(sharpe),         (sharpe > 1) if not math.isnan(sharpe) else None),
        ("Max Drawdown",  f"{max_dd:.1f}%" if not math.isnan(max_dd) else "N/A",
                          False if not math.isnan(max_dd) else None),
    ]
    _tw = 3.0
    for _i, (_lbl, _val, _pos) in enumerate(_bl_kpis):
        _x = 0.3 + _i * (_tw + 0.18)
        _rect(sl, _x, 4.05, _tw, 1.55, fill_rgb=C_WHITE,
              line_rgb=RGBColor(0xE2, 0xE8, 0xF0), line_width=1)
        _rect(sl, _x, 4.05, _tw, 0.07, fill_rgb=C_ACCENT)
        _vcol = C_GREEN if _pos is True else (C_RED if _pos is False else C_NAVY)
        _text_box(sl, _lbl.upper(), _x + 0.2, 4.28, _tw - 0.35, 0.3,
                  font_size=9.5, bold=True, color=C_GREY_TEXT)
        _text_box(sl, _val, _x + 0.2, 4.68, _tw - 0.35, 0.75,
                  font_size=26, bold=True, color=_vcol)

    # ── Valuation — the conclusion slide, straight after the executive summary ─
    if dcf is not None:
        _build_valuation_slide(prs, ticker, dcf, _next_page(), total_slides)

    # ── Fundamentals & Valuation (only when EDGAR/Polygon data is available) ───
    if fundamentals and fundamentals.get("ok"):
        _build_fundamentals_slide(prs, ticker, fundamentals, _next_page(), total_slides)

    # ── Slide 2: Company Snapshot ─────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Company Snapshot", ticker)
    _slide_footer(sl, _next_page(), total_slides)

    desc = cd.get("Description", "No description available.")
    if len(desc) > 900:
        desc = desc[:900].rsplit(" ", 1)[0] + "…"

    # Full-height About panel on the left + key-facts table on the right. The
    # executive verdict lives on the Bottom Line slide, so it isn't repeated here.
    _rect(sl, 0.3, 1.3, 8.5, 5.55, fill_rgb=C_LIGHT)
    _text_box(sl, "About", 0.5, 1.42, 8, 0.3, font_size=11, bold=True, color=C_NAVY)
    _text_box(sl, desc, 0.5, 1.85, 8.1, 4.85, font_size=10.5, color=C_DARK_TEXT)

    mc_raw = cd.get("Market Cap")
    if isinstance(mc_raw, (int, float)) and mc_raw >= 1e12:
        mc_str = f"${mc_raw/1e12:.2f}T"
    elif isinstance(mc_raw, (int, float)) and mc_raw >= 1e9:
        mc_str = f"${mc_raw/1e9:.1f}B"
    elif isinstance(mc_raw, (int, float)):
        mc_str = f"${mc_raw/1e6:.0f}M"
    else:
        mc_str = "N/A"
    emp_raw = cd.get("Employees")
    emp_str = f"{int(emp_raw):,}" if isinstance(emp_raw, (int, float)) else "N/A"
    country = cd.get("Country", "N/A")
    if isinstance(country, str) and len(country) <= 3:
        country = country.upper()   # "us" -> "US"

    info_pairs = [
        ("Ticker",    ticker),
        ("Company",   cd.get("Name",     "N/A")),
        ("Sector",    cd.get("Sector",   "N/A")),
        ("Exchange",  cd.get("Exchange", "N/A")),
        ("Market Cap", mc_str),
        ("Employees", emp_str),
        ("Country",   country),
        ("Website",   cd.get("Website",  "N/A")),
    ]
    _kv_block(sl, info_pairs, 9.15, 1.3, 3.85, col_w=1.6, row_h=0.69)

    # ── Slide 3: Key Metrics ──────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Key Performance Metrics", f"{ticker}  ·  {period_label}")
    _slide_footer(sl, _next_page(), total_slides)

    rsi_val = latest.get("RSI14")
    try:
        rsi_val = float(rsi_val) if rsi_val is not None and not (isinstance(rsi_val, float) and math.isnan(rsi_val)) else None
    except Exception:
        rsi_val = None

    rsi_str = f"{rsi_val:.1f}" if rsi_val is not None else "N/A"
    rsi_pos = (rsi_val < 70) if rsi_val is not None else None

    col1 = [
        ("Current Price",       f"${latest['Close']:,.2f}"),
        ("Period Return",       _fmt_pct(period_ret),    period_ret >= 0),
        ("52-Week High",        f"${latest.get('52W_High', 0):,.2f}" if latest.get('52W_High') else "N/A"),
        ("52-Week Low",         f"${latest.get('52W_Low', 0):,.2f}"  if latest.get('52W_Low')  else "N/A"),
        ("% From 52W High",     _fmt_pct((latest.get("Pct_From_52W_High") or 0) * 100),
                                (latest.get("Pct_From_52W_High") or -1) > -0.05),
    ]
    col2 = [
        ("Sharpe Ratio",        _fmt_ratio(sharpe),  sharpe > 1 if not math.isnan(sharpe) else None),
        ("Sortino Ratio",       _fmt_ratio(sortino), sortino > 1 if not math.isnan(sortino) else None),
        ("Ann. Return",         _fmt_pct(ann_ret * 100), ann_ret >= 0),
        ("Ann. Volatility",     _fmt_pct(ann_std * 100)),
        ("Max Drawdown (60d)",  _fmt_pct(max_dd),    max_dd > -20),
    ]
    col3 = [
        ("20-Day MA",   f"${latest.get('MA20', 0):,.2f}"  if latest.get('MA20')  else "N/A"),
        ("50-Day MA",   f"${latest.get('MA50', 0):,.2f}"  if latest.get('MA50')  else "N/A"),
        ("200-Day MA",  f"${latest.get('MA200', 0):,.2f}" if latest.get('MA200') else "N/A"),
        ("RSI (14)",    rsi_str, rsi_pos),
        ("BB %B",       f"{latest.get('BB_Pct', 0):.2f}"  if latest.get('BB_Pct') is not None else "N/A"),
    ]

    _text_box(sl, "Price & Returns",   0.3,  1.5, 4.1, 0.28, font_size=11, bold=True, color=C_NAVY)
    _text_box(sl, "Risk Metrics",      4.65, 1.5, 4.1, 0.28, font_size=11, bold=True, color=C_NAVY)
    _text_box(sl, "Technical Levels",  9.0,  1.5, 4.1, 0.28, font_size=11, bold=True, color=C_NAVY)

    # Taller rows so the three columns fill the slide instead of clustering at the top.
    _kv_block(sl, col1, 0.3,  1.9, 4.2, col_w=1.9, row_h=0.92)
    _kv_block(sl, col2, 4.65, 1.9, 4.2, col_w=1.9, row_h=0.92)
    _kv_block(sl, col3, 9.0,  1.9, 4.2, col_w=1.9, row_h=0.92)

    # ── Slide 4: Price & Moving Averages ──────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Price & Moving Averages", f"{ticker}  ·  {period_label}")
    _slide_footer(sl, _next_page(), total_slides)

    buf = _price_ma_chart(df, ticker, w=12, h=5.0)
    _add_image(sl, buf, 0.6, 1.25, 12.0, 5.0)

    # ── Slide 5: Volume ───────────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Volume Analysis", f"{ticker}  ·  {period_label}")
    _slide_footer(sl, _next_page(), total_slides)

    buf_vol = _volume_chart(df, ticker, w=12, h=3.2)
    _add_image(sl, buf_vol, 0.6, 1.25, 12.0, 3.2)

    buf_dd = _drawdown_chart(df, ticker, w=12, h=2.5)
    if buf_dd:
        _add_image(sl, buf_dd, 0.6, 4.55, 12.0, 2.5)

    # ── Slide 6: Bollinger Bands ──────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Bollinger Bands  ·  RSI", f"{ticker}  ·  {period_label}")
    _slide_footer(sl, _next_page(), total_slides)

    buf_bb = _bollinger_chart(df, ticker, w=12, h=3.5)
    if buf_bb:
        _add_image(sl, buf_bb, 0.6, 1.25, 12.0, 3.5)

    buf_rsi = _rsi_chart(df, ticker, w=12, h=2.4)
    if buf_rsi:
        _add_image(sl, buf_rsi, 0.6, 4.8, 12.0, 2.4)

    # ── Slide 7: Cumulative Return vs Benchmarks ──────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Cumulative Return vs Benchmarks", f"{ticker}  ·  {period_label}")
    _slide_footer(sl, _next_page(), total_slides)

    buf_cum = _cumulative_chart(df, ticker, w=12, h=5.0)
    if buf_cum:
        _add_image(sl, buf_cum, 0.6, 1.25, 12.0, 5.0)

    # ── Slide 8: Monte Carlo Forecast ─────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Monte Carlo Forecast", f"{ticker}  ·  1,000 Simulations")
    _slide_footer(sl, _next_page(), total_slides)

    if mc_sim_df is not None and mc_summary:
        buf_mc = _monte_carlo_chart(mc_sim_df, mc_summary, ticker, w=8.5, h=5.0)
        if buf_mc:
            _add_image(sl, buf_mc, 0.6, 1.25, 8.5, 5.0)

        mc_pairs = [
            ("Last Price",      mc_summary.get("Last Price", "N/A")),
            ("Median (P50)",    mc_summary.get("Median (P50)", "N/A")),
            ("Bear (P5)",       mc_summary.get("Bear Case (P5)", "N/A")),
            ("Bull (P75)",      mc_summary.get("Bull Case (P75)", "N/A")),
            ("Best (P95)",      mc_summary.get("Best Case (P95)", "N/A")),
            ("Prob. of Gain",   mc_summary.get("Prob. of Gain", "N/A")),
            ("Ann. Volatility", mc_summary.get("Ann. Volatility", "N/A")),
        ]
        _text_box(sl, "Scenario Summary", 9.35, 1.25, 3.6, 0.3,
                  font_size=10, bold=True, color=C_NAVY)
        _kv_block(sl, mc_pairs, 9.35, 1.6, 3.6, col_w=1.4, row_h=0.52)
    else:
        _text_box(sl, "Monte Carlo simulation was not run for this analysis.\n\n"
                  "Enable it in the sidebar options to generate probabilistic forecasts.",
                  0.6, 2.5, 12, 1.5, font_size=13, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    # ── Slide 9: News Headlines ───────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Recent News Headlines", ticker)
    _slide_footer(sl, _next_page(), total_slides)

    if news_list:
        shown = news_list[:8]
        for i, item in enumerate(shown):
            row_t = 1.3 + i * 0.7
            bg = C_LIGHT if i % 2 == 0 else C_WHITE
            _rect(sl, 0.3, row_t, 12.7, 0.66, fill_rgb=bg)
            date_str = item.get("Date", "")[:10]
            pub      = item.get("Publisher", "")
            headline = item.get("Headline", "")
            if len(headline) > 110:
                headline = headline[:110] + "…"
            _text_box(sl, f"{date_str}  ·  {pub}", 0.45, row_t + 0.04, 3.5, 0.26,
                      font_size=8, color=C_GREY_TEXT)
            _text_box(sl, headline, 0.45, row_t + 0.3, 12.3, 0.32,
                      font_size=9.5, bold=True, color=C_DARK_TEXT)
    else:
        _text_box(sl, "No news headlines available for this ticker.",
                  0.6, 3.0, 12, 0.5, font_size=13, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    # ── Slide 10: Disclaimer ──────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x0F, 0x17, 0x2A))
    _rect(sl, 0, 0, 13.33, 0.06, fill_rgb=C_ACCENT)
    _rect(sl, 0, 7.44, 13.33, 0.06, fill_rgb=C_ACCENT)

    _text_box(sl, "QUANTWIZARD", 0.6, 0.7, 12, 0.4,
              font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _text_box(sl, "Important Disclaimer", 0.6, 1.2, 12, 0.55,
              font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    disclaimer = (
        "This report has been generated by QuantWizard for informational and educational purposes "
        "only. It does not constitute financial, investment, legal, or tax advice. The information "
        "presented is derived from third-party data sources (Polygon, Yahoo Finance, Finnhub, and "
        "SEC EDGAR) and is believed to be accurate but is not guaranteed.\n\n"
        "Past performance is not indicative of future results. All investments involve risk, "
        "including the possible loss of principal. You should not make any investment decision "
        "based solely on the information in this report.\n\n"
        "QuantWizard is not a registered investment adviser, broker-dealer, or financial planner. "
        "Always consult a qualified financial professional before making investment decisions."
    )
    _text_box(sl, disclaimer, 1.0, 2.0, 11.3, 4.2,
              font_size=10, color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.LEFT)

    _text_box(sl, f"© {datetime.now().year} QuantWizard  ·  quantwizard.app",
              0.6, 6.6, 12, 0.3, font_size=9, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Portfolio Deck ────────────────────────────────────────────────────────────

def _risk_label_from_tolerance(tol):
    """Map the app's 1-10 risk_tolerance to the same labels the UI shows."""
    try:
        tol = int(tol)
    except (TypeError, ValueError):
        return "Moderate"
    if tol <= 3:
        return "Conservative"
    if tol <= 6:
        return "Moderate"
    if tol <= 9:
        return "Aggressive"
    return "Ultra Aggressive"


def _horizon_to_years(horizon, default=5):
    """The app stores horizon as a string ('5 years', '20+ years'); pull the number."""
    if isinstance(horizon, (int, float)):
        return int(horizon)
    if isinstance(horizon, str):
        m = re.search(r"\d+", horizon)
        if m:
            return int(m.group())
    return default


def build_portfolio_pptx(preferences, final_weights, stock_metrics,
                          backtest_df=None, backtest_metrics=None,
                          mc_sim_df=None, mc_summary=None, milestones=None,
                          corr_matrix=None, diversification_score=None,
                          ticker_info=None):
    """Build a professional portfolio analysis PowerPoint. Returns BytesIO."""
    if not PPTX_AVAILABLE:
        raise RuntimeError("python-pptx is not installed.")

    prs   = _new_prs()
    prefs = preferences or {}
    ti    = ticker_info or {}
    bm    = backtest_metrics or {}
    tickers = list(final_weights.keys()) if final_weights else []

    total_slides = 9
    page = [1]   # cover is page 1; next slide starts at 2

    def _next_page():
        page[0] += 1
        return page[0]

    # Read the keys the app actually sets (starting_capital / risk_tolerance /
    # horizon), falling back to the legacy names then sane defaults. Previously the
    # deck read investment_amount/risk_label/horizon_years — which the app never
    # sets — so every report silently showed $10,000 / Moderate / 5yr.
    inv_amount  = prefs.get("starting_capital", prefs.get("investment_amount", 10000))
    risk_label  = prefs.get("risk_label") or _risk_label_from_tolerance(prefs.get("risk_tolerance", 5))
    horizon_yrs = _horizon_to_years(prefs.get("horizon", prefs.get("horizon_years")), 5)

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill_rgb=C_NAVY)
    _rect(sl, 0, 5.5, 13.33, 2.0, fill_rgb=RGBColor(0x0B, 0x14, 0x24))   # KPI band
    _rect(sl, 0, 5.5, 13.33, 0.03, fill_rgb=C_ACCENT)
    _rect(sl, 0.6, 1.3, 0.09, 2.5, fill_rgb=C_ACCENT)

    _text_box(sl, "QUANTWIZARD", 0.9, 0.85, 12, 0.5,
              font_size=14, bold=True, color=C_ACCENT)
    _text_box(sl, "Portfolio Analysis", 0.86, 1.3, 12, 1.0,
              font_size=52, bold=True, color=C_WHITE)
    _text_box(sl, "Optimized Portfolio Report", 0.9, 2.45, 12, 0.45,
              font_size=16, italic=True, color=C_ACCENT)

    holdings_str = "   ·   ".join(tickers[:12])
    _text_box(sl, holdings_str, 0.9, 3.05, 12, 0.5,
              font_size=11, color=RGBColor(0xB0, 0xC4, 0xDE))

    meta_parts = [
        f"Risk Profile: {risk_label}",
        f"Horizon: {horizon_yrs}yr",
        f"Amount: ${inv_amount:,.0f}",
        f"{len(tickers)} Holdings",
    ]
    _text_box(sl, "   ·   ".join(meta_parts), 0.9, 3.75, 12, 0.3,
              font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))

    _text_box(sl, f"Generated {datetime.now().strftime('%B %d, %Y')}   ·   "
                  f"Multi-source data: Polygon · Yahoo Finance · Finnhub · SEC EDGAR",
              0.9, 4.3, 12, 0.3, font_size=9.5, color=C_GREY_TEXT)
    _text_box(sl, "For informational purposes only. Not financial advice.",
              0.9, 4.63, 12, 0.3, font_size=8.5, italic=True, color=C_GREY_TEXT)

    # KPI stat band from the backtest metrics (falls back to portfolio facts)
    def _bvf(k, d=0.0):
        try:    return float(bm.get(k, d))
        except Exception: return d
    if bm:
        _tr, _dd = _bvf("Total Return"), _bvf("Max Drawdown")
        _pk = [
            ("TOTAL RETURN", f"{_tr:+.1f}%", C_GREEN if _tr >= 0 else C_RED),
            ("CAGR",         f"{_bvf('Ann. Return'):+.1f}%", C_WHITE),
            ("SHARPE RATIO", f"{_bvf('Sharpe Ratio'):.2f}", C_WHITE),
            ("MAX DRAWDOWN", f"{_dd:.1f}%", C_RED),
        ]
    else:
        _pk = [
            ("HOLDINGS",        str(len(tickers)),  C_WHITE),
            ("RISK PROFILE",    risk_label,          C_WHITE),
            ("HORIZON",         f"{horizon_yrs}yr",  C_WHITE),
            ("DIVERSIFICATION", f"{diversification_score}/10" if diversification_score else "N/A", C_WHITE),
        ]
    for _i, (_l, _v, _c) in enumerate(_pk):
        _x = 0.9 + _i * 3.05
        _text_box(sl, _v, _x, 5.85, 2.9, 0.75, font_size=30, bold=True, color=_c)
        _text_box(sl, _l, _x + 0.03, 6.65, 2.9, 0.3, font_size=9.5, bold=True, color=C_ACCENT)

    _text_box(sl, "1", 12.9, 7.15, 0.4, 0.25, font_size=8, color=C_GREY_TEXT, align=PP_ALIGN.RIGHT)

    # ── Slide 2: Portfolio Allocation ─────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Portfolio Allocation", f"{len(tickers)} Holdings  ·  {risk_label} Risk")
    _slide_footer(sl, _next_page(), total_slides)

    if final_weights:
        buf_pie = _alloc_pie_chart(final_weights, ti, w=6.5, h=5.0)
        _add_image(sl, buf_pie, 0.3, 1.25, 6.5, 5.0)

    alloc_pairs = [(tk, f"{wt*100:.1f}%") for tk, wt in
                   sorted(final_weights.items(), key=lambda x: -x[1])]
    _text_box(sl, "Weights", 7.2, 1.25, 5.8, 0.28, font_size=10, bold=True, color=C_NAVY)
    _kv_block(sl, alloc_pairs, 7.2, 1.6, 5.8, col_w=2.5,
              row_h=min(0.52, 5.0 / max(len(alloc_pairs), 1)))

    # ── Slide 3: Holdings Breakdown ───────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Holdings Breakdown", "Individual Performance Metrics")
    _slide_footer(sl, _next_page(), total_slides)

    buf_bar = _holdings_bar_chart(stock_metrics, w=12, h=3.8)
    if buf_bar:
        _add_image(sl, buf_bar, 0.6, 1.25, 12.0, 3.8)

    if stock_metrics:
        metrics_pairs = []
        for tk in tickers[:8]:
            m = stock_metrics.get(tk, {})
            ann_r = m.get("ann_return", 0)   # already a percentage
            sharpe_v = m.get("sharpe", float("nan"))
            sharpe_s = f"{sharpe_v:.2f}" if not (isinstance(sharpe_v, float) and math.isnan(sharpe_v)) else "N/A"
            metrics_pairs.append((tk, f"Ret: {ann_r:+.1f}%  ·  Sharpe: {sharpe_s}",
                                   ann_r >= 0))
        _text_box(sl, "Per-Holding Summary", 0.3, 5.2, 12, 0.28,
                  font_size=10, bold=True, color=C_NAVY)
        _kv_block(sl, metrics_pairs, 0.3, 5.5, 12.7, col_w=1.0, row_h=0.3)

    # ── Slide 4: Portfolio Metrics ────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Portfolio Performance Metrics", "Backtest Results")
    _slide_footer(sl, _next_page(), total_slides)

    def _bv(key, default="N/A"):
        v = bm.get(key, default)
        return v if v is not None else default

    def _bvp(key, signed=True):
        """Format a backtest metric (stored as a bare number) as a percent."""
        v = bm.get(key)
        if v is None or v == "N/A":
            return "N/A"
        try:
            f = float(v)
        except (TypeError, ValueError):
            return str(v)
        return f"{f:+.1f}%" if signed else f"{f:.1f}%"

    col1 = [
        ("Total Return",         _bvp("Total Return")),
        ("Ann. Return",          _bvp("Ann. Return")),
        ("vs S&P 500",           _bvp("vs S&P 500")),
        ("Final Portfolio Value", f"${float(str(_bv('Final Value')).replace('$','').replace(',','').replace('%','')):,.0f}"
                                  if _bv("Final Value") != "N/A" else "N/A"),
    ]
    col2 = [
        ("Sharpe Ratio",         _bv("Sharpe Ratio")),
        ("Max Drawdown",         _bvp("Max Drawdown")),
        ("Best Month",           _bvp("Best Month")),
        ("% Months Positive",    _bvp("% Months Positive", signed=False)),
    ]
    col3 = [
        ("Risk Profile",         risk_label),
        ("Horizon",              f"{horizon_yrs} years"),
        ("Investment Amount",    f"${inv_amount:,.0f}"),
        ("Diversification",      f"{diversification_score:.1f}/10" if diversification_score else "N/A"),
    ]

    _text_box(sl, "Return Metrics",    0.3,  1.5, 4.2, 0.28, font_size=11, bold=True, color=C_NAVY)
    _text_box(sl, "Risk Metrics",      4.75, 1.5, 4.2, 0.28, font_size=11, bold=True, color=C_NAVY)
    _text_box(sl, "Portfolio Profile", 9.2,  1.5, 4.2, 0.28, font_size=11, bold=True, color=C_NAVY)
    # Taller rows so the three columns fill the slide instead of clustering at top.
    _kv_block(sl, col1, 0.3,  1.9, 4.3, col_w=2.0, row_h=0.92)
    _kv_block(sl, col2, 4.75, 1.9, 4.3, col_w=2.0, row_h=0.92)
    _kv_block(sl, col3, 9.2,  1.9, 4.1, col_w=2.0, row_h=0.92)

    # ── Slide 5: Backtest Performance Chart ───────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Backtest — Portfolio Growth", "Account value vs S&P 500 (same contributions)")
    _slide_footer(sl, _next_page(), total_slides)

    buf_bt = _portfolio_performance_chart(backtest_df, w=12, h=5.0)
    if buf_bt:
        _add_image(sl, buf_bt, 0.6, 1.25, 12.0, 5.0)
    else:
        _text_box(sl, "Backtest chart unavailable.", 0.6, 3.5, 12, 0.5,
                  font_size=13, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    # ── Slide 6: Monte Carlo ──────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Monte Carlo Portfolio Forecast", "Probabilistic Return Scenarios")
    _slide_footer(sl, _next_page(), total_slides)

    if mc_sim_df is not None:
        buf_mc = _monte_carlo_chart(mc_sim_df, mc_summary or {}, "Portfolio", w=8.5, h=5.0)
        if buf_mc:
            _add_image(sl, buf_mc, 0.6, 1.25, 8.5, 5.0)

        if milestones:
            # Each milestone is a dict (P5/P25/P50/.../prob_gain); show the median
            # projected value per horizon rather than dumping the raw dict.
            def _num(x):
                try:
                    return float(str(x).replace("$", "").replace(",", ""))
                except (TypeError, ValueError):
                    return None
            ms_pairs = []
            for k, v in milestones.items():
                if isinstance(v, dict):
                    med = _num(v.get("P50"))
                    pg  = str(v.get("prob_gain", "")).strip()
                    val = (f"${med:,.0f}" + (f"  ·  {pg} gain" if pg else "")) if med is not None else "—"
                else:
                    val = str(v)
                ms_pairs.append((k.replace("yr", "-Year"), val))
            _text_box(sl, "Median Projected Value", 9.35, 1.25, 3.6, 0.3,
                      font_size=10, bold=True, color=C_NAVY)
            _kv_block(sl, ms_pairs, 9.35, 1.6, 3.6, col_w=1.3, row_h=0.62)
    else:
        _text_box(sl, "Monte Carlo simulation not available.",
                  0.6, 3.5, 12, 0.5, font_size=13, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    # ── Slide 7: Correlation Matrix ───────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Correlation Matrix", "Daily Return Correlation Between Holdings")
    _slide_footer(sl, _next_page(), total_slides)

    if corr_matrix is not None and not corr_matrix.empty:
        labels = list(corr_matrix.columns)
        n      = len(labels)
        fig, ax = plt.subplots(figsize=(10, 5))
        data    = corr_matrix.values.astype(float)
        im      = ax.imshow(data, cmap="RdYlGn", vmin=-1, vmax=1, aspect="auto")
        ax.set_xticks(range(n)); ax.set_xticklabels(labels, rotation=40, ha="right", fontsize=8)
        ax.set_yticks(range(n)); ax.set_yticklabels(labels, fontsize=8)
        for i in range(n):
            for j in range(n):
                val = data[i, j]
                ax.text(j, i, f"{val:.2f}", ha="center", va="center",
                        fontsize=7.5, color="black" if abs(val) < 0.7 else "white",
                        fontweight="bold")
        plt.colorbar(im, ax=ax, fraction=0.03, pad=0.02)
        ax.set_title("Correlation Matrix (Daily Returns)", fontsize=11,
                     fontweight="bold", color="#1F4E79", pad=8)
        fig.patch.set_facecolor("#FFFFFF")
        fig.tight_layout()
        buf_corr = _fig_to_buf(fig)
        _add_image(sl, buf_corr, 1.0, 1.3, 11.3, 5.5)
    else:
        _text_box(sl, "Correlation data unavailable.", 0.6, 3.5, 12, 0.5,
                  font_size=13, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    # ── Slide 8: Rebalancing & Notes ──────────────────────────────────────────
    sl = _blank_slide(prs)
    _slide_header(sl, "Rebalancing Guide", "Suggested Target Weights & Notes")
    _slide_footer(sl, _next_page(), total_slides)

    _rect(sl, 0.3, 1.3, 12.7, 0.38, fill_rgb=C_NAVY)
    _text_box(sl, "Target Allocation",    0.5,  1.35, 3.0, 0.3, font_size=9, bold=True, color=C_WHITE)
    _text_box(sl, "Ticker",               3.6,  1.35, 1.5, 0.3, font_size=9, bold=True, color=C_WHITE)
    _text_box(sl, "Weight",               5.2,  1.35, 1.5, 0.3, font_size=9, bold=True, color=C_WHITE)
    _text_box(sl, "$ Amount",             6.8,  1.35, 2.0, 0.3, font_size=9, bold=True, color=C_WHITE)
    _text_box(sl, "Ann. Return",          8.9,  1.35, 1.8, 0.3, font_size=9, bold=True, color=C_WHITE)
    _text_box(sl, "Sharpe",              10.8,  1.35, 1.5, 0.3, font_size=9, bold=True, color=C_WHITE)

    for i, (tk, wt) in enumerate(sorted(final_weights.items(), key=lambda x: -x[1])):
        row_t  = 1.72 + i * 0.43
        bg     = C_LIGHT if i % 2 == 0 else C_WHITE
        _rect(sl, 0.3, row_t, 12.7, 0.41, fill_rgb=bg)
        m      = stock_metrics.get(tk, {})
        name   = ti.get(tk, {}).get("name", "")
        if len(name) > 28:
            name = name[:28] + "…"
        ann_r  = m.get("ann_return", 0)   # already a percentage
        sh_v   = m.get("sharpe", float("nan"))
        sh_s   = f"{sh_v:.2f}" if not (isinstance(sh_v, float) and math.isnan(sh_v)) else "N/A"
        _text_box(sl, name,                0.45, row_t + 0.06, 3.1, 0.3, font_size=9, color=C_DARK_TEXT)
        _text_box(sl, tk,                  3.6,  row_t + 0.06, 1.5, 0.3, font_size=9, bold=True, color=C_NAVY)
        _text_box(sl, f"{wt*100:.1f}%",    5.2,  row_t + 0.06, 1.5, 0.3, font_size=9, color=C_DARK_TEXT)
        _text_box(sl, f"${wt*inv_amount:,.0f}", 6.8, row_t + 0.06, 2.0, 0.3, font_size=9, color=C_DARK_TEXT)
        clr = C_GREEN if ann_r >= 0 else C_RED
        _text_box(sl, f"{ann_r:+.1f}%",   8.9,  row_t + 0.06, 1.8, 0.3, font_size=9, bold=True, color=clr)
        _text_box(sl, sh_s,               10.8,  row_t + 0.06, 1.5, 0.3, font_size=9, color=C_DARK_TEXT)

    note = ("Target weights reflect the optimized portfolio based on your risk profile and historical data. "
            "Rebalance when individual positions drift more than 5% from target weights. "
            "Past performance does not guarantee future results.")
    _rect(sl, 0.3, 6.55, 12.7, 0.6, fill_rgb=RGBColor(0xFF, 0xFB, 0xEB))
    _text_box(sl, f"⚠  {note}", 0.5, 6.6, 12.3, 0.5, font_size=8.5,
              italic=True, color=RGBColor(0x92, 0x40, 0x0E))

    # ── Slide 9: Disclaimer ───────────────────────────────────────────────────
    sl = _blank_slide(prs)
    _rect(sl, 0, 0, 13.33, 7.5, fill_rgb=RGBColor(0x0F, 0x17, 0x2A))
    _rect(sl, 0, 0,    13.33, 0.06, fill_rgb=C_ACCENT)
    _rect(sl, 0, 7.44, 13.33, 0.06, fill_rgb=C_ACCENT)

    _text_box(sl, "QUANTWIZARD", 0.6, 0.7, 12, 0.4,
              font_size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _text_box(sl, "Important Disclaimer", 0.6, 1.2, 12, 0.55,
              font_size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    disclaimer = (
        "This portfolio analysis has been generated by QuantWizard for informational and "
        "educational purposes only. It does not constitute financial, investment, legal, or "
        "tax advice.\n\n"
        "Optimized portfolio weights are based on historical data and mathematical models. "
        "They are not guarantees of future performance. All investments involve risk, "
        "including the possible loss of principal.\n\n"
        "The rebalancing suggestions in this report are informational targets based on "
        "your stated preferences — not personalized investment advice. Always consult a "
        "qualified financial professional before making investment decisions.\n\n"
        "QuantWizard is not a registered investment adviser or broker-dealer."
    )
    _text_box(sl, disclaimer, 1.0, 2.0, 11.3, 4.5,
              font_size=10, color=RGBColor(0xB0, 0xC4, 0xDE))

    _text_box(sl, f"© {datetime.now().year} QuantWizard  ·  quantwizard.app",
              0.6, 6.6, 12, 0.3, font_size=9, color=C_GREY_TEXT, align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

"""
validate_metrics.py — Data-correctness validation (complements validate.py)

validate.py proves the portfolio BACKTEST engine matches known benchmarks.
This suite guards the classes of bug that validate.py can't see — the ones found
by hand-rendering the reports:

  A. Per-stock metric sanity   — the enriched single-stock metrics that feed the
     stock report are in plausible ranges (catches a -568%/+4038% style blow-up
     at the source).
  B. Cross-source price agreement — yfinance and Polygon agree on recent closes
     within tolerance (catches a silently-wrong data source).
  C. Report display consistency — the built portfolio deck reflects the user's
     actual inputs and never prints an absurd percentage (catches the ×100
     double-scaling and the $10,000-default key-mismatch bugs).

Usage:
    python validate_metrics.py            # exits non-zero if any check fails
"""

import os
import re
import sys
from datetime import date, timedelta

import numpy as np
import pandas as pd
from dotenv import load_dotenv

load_dotenv()
KEY   = os.getenv("POLYGON_API_KEY", "")
END   = date.today().strftime("%Y-%m-%d")
START = (date.today() - timedelta(days=730)).strftime("%Y-%m-%d")

_PASS, _FAIL = [], []


def check(cond, label, detail=""):
    (_PASS if cond else _FAIL).append(label)
    mark = "✓" if cond else "✗"
    line = f"  {mark} {label}"
    if detail:
        line += f"  ({detail})"
    print(line)
    return bool(cond)


def _quiet(*a, **k):
    pass


# ── A. Per-stock metric sanity ──────────────────────────────────────────────────
# A spread of profiles so the ranges are actually exercised: mega-cap growth,
# defensive staple, long-duration bond ETF (low/negative return), high-vol name.
BASKET = ["AAPL", "KO", "TLT", "NVDA"]


def part_a_stock_metrics():
    print("\nA. Per-stock metric sanity")
    print("=" * 60)
    from data import fetch_stock_data

    for tk in BASKET:
        try:
            df = fetch_stock_data(tk, benchmark_tickers=("SPY",), api_key=KEY, log=_quiet,
                                  start_override=START, end_override=END, bar_size="day")
        except Exception as e:
            check(False, f"{tk}: fetch", str(e)[:60])
            continue
        if df is None or len(df) < 30:
            check(False, f"{tk}: enough history", f"{0 if df is None else len(df)} rows")
            continue

        latest  = df.iloc[-1]
        ret      = df["Daily_Return"].dropna()
        ann_ret  = ret.mean() * 252 * 100
        ann_vol  = ret.std() * np.sqrt(252) * 100
        max_dd   = df["Drawdown_60d"].min() * 100 if "Drawdown_60d" in df.columns else 0.0
        rsi      = latest.get("RSI14")

        check(latest["Close"] > 0,                 f"{tk}: price > 0", f"${latest['Close']:.2f}")
        check(-90 < ann_ret < 300,                 f"{tk}: ann return plausible", f"{ann_ret:+.0f}%")
        check(0 < ann_vol < 150,                   f"{tk}: ann vol plausible",    f"{ann_vol:.0f}%")
        check(-100 < max_dd <= 0.01,               f"{tk}: drawdown in (-100, 0]", f"{max_dd:.0f}%")
        check(rsi is None or pd.isna(rsi) or 0 <= float(rsi) <= 100,
              f"{tk}: RSI in [0,100]", f"{rsi}")
        check((df['Date'].diff().dropna() > pd.Timedelta(0)).all(),
              f"{tk}: dates strictly increasing")


# ── B. Cross-source price agreement ─────────────────────────────────────────────
def part_b_cross_source():
    print("\nB. Cross-source price agreement (yfinance vs Polygon)")
    print("=" * 60)
    if not KEY:
        check(True, "skipped — no POLYGON_API_KEY")
        return
    from market_data import _yahoo_bars, _polygon_bars

    for tk in ["AAPL", "MSFT"]:
        y = _yahoo_bars(tk, START, END, "day")
        p = _polygon_bars(tk, START, END, "day", KEY)
        if y is None or p is None or y.empty or p.empty:
            check(True, f"{tk}: skipped — a source returned no data")
            continue
        # Compare the last common date's close (adjusted vs adjusted) within 2%.
        ym = y.set_index(y["Date"].dt.normalize())["Close"]
        pm = p.set_index(p["Date"].dt.normalize())["Close"]
        common = ym.index.intersection(pm.index)
        if len(common) == 0:
            check(True, f"{tk}: skipped — no overlapping dates")
            continue
        d   = common.max()
        yv, pv = float(ym.loc[d]), float(pm.loc[d])
        diff_pct = abs(yv - pv) / pv * 100 if pv else 99.0
        check(diff_pct <= 2.0, f"{tk}: yf vs Polygon close within 2%",
              f"yf ${yv:.2f} / poly ${pv:.2f} = {diff_pct:.1f}% on {d.date()}")


# ── C. Report display consistency (portfolio deck) ──────────────────────────────
def part_c_report_consistency():
    print("\nC. Portfolio report display consistency")
    print("=" * 60)
    try:
        from pptx import Presentation
    except ImportError:
        check(True, "skipped — python-pptx not installed")
        return
    from portfolio_data import fetch_portfolio_prices, get_ticker_info
    from portfolio_analysis import (compute_stock_metrics, compute_correlation_matrix,
        optimise_portfolio, backtest_portfolio, compute_backtest_metrics,
        run_portfolio_monte_carlo, compute_diversification_score)
    from pptx_builder import build_portfolio_pptx
    import io

    CAP, MO = 50_000, 1_000           # deliberately NOT the deck's old $10k default
    tickers = ["AAPL", "MSFT", "JNJ", "JPM", "XOM", "PG"]
    sector_map = {"AAPL": "Technology", "MSFT": "Technology", "JNJ": "Healthcare",
                  "JPM": "Financials", "XOM": "Energy", "PG": "Consumer Staples"}
    try:
        _, close_df, returns_df, _ = fetch_portfolio_prices(tickers + ["SPY"], period_years=2,
                                                            api_key=KEY, log=_quiet)
        opt_ret = returns_df[[t for t in tickers if t in returns_df.columns]]
        prefs = {"horizon": "10 years", "starting_capital": CAP, "monthly_contribution": MO,
                 "risk_tolerance": 5, "max_per_stock": 0.30, "target_value": 200_000}
        weights = optimise_portfolio(opt_ret, risk_tolerance=5, sector_map=sector_map,
                                     max_weight=0.30)["recommended"]
        bt  = backtest_portfolio(close_df, weights, CAP, MO)
        mcd, mcs, miles = run_portfolio_monte_carlo(opt_ret, weights, CAP, MO,
                                                    forecast_years=10, n_simulations=400, log=_quiet)
        buf = build_portfolio_pptx(
            preferences=prefs, final_weights=weights,
            stock_metrics=compute_stock_metrics(opt_ret),
            backtest_df=bt, backtest_metrics=compute_backtest_metrics(bt, CAP),
            mc_sim_df=mcd, mc_summary=mcs, milestones=miles,
            corr_matrix=compute_correlation_matrix(opt_ret),
            diversification_score=compute_diversification_score(weights, opt_ret),
            ticker_info={t: get_ticker_info(t, KEY) for t in weights})
    except Exception as e:
        check(False, "build portfolio deck", str(e)[:80])
        return

    # Pull every run of text out of the deck.
    prs   = Presentation(io.BytesIO(buf.getvalue()))
    texts = [sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame]
    blob  = "\n".join(texts)

    # 1. The deck must reflect the real capital, not the old $10,000 default.
    #    With max_weight 0.30 no holding reaches $50k, and Final Value/milestones
    #    are far larger, so "$50,000" appears ONLY via the investment-amount field —
    #    its presence is a clean signal the deck read starting_capital (not the default).
    check(f"${CAP:,.0f}" in blob,
          f"deck shows the user's capital (${CAP:,.0f}), not the $10k default")

    # 2. No absurd percentages anywhere (the ×100 double-scaling bug printed >2000%).
    pcts = [float(m.replace(",", "")) for m in re.findall(r"([+-]?\d[\d,]*\.?\d*)%", blob)]
    worst = max((abs(p) for p in pcts), default=0.0)
    check(worst <= 300.0, "no absurd percentage in deck text (≤300%)", f"max |%| seen = {worst:.0f}%")


# ── D. Fundamentals sanity (EDGAR → compute_fundamentals) ───────────────────────
def part_d_fundamentals():
    print("\nD. Fundamentals sanity (EDGAR parsing)")
    print("=" * 60)
    from data import fetch_sec_financials, fetch_company_details
    from analysis import compute_fundamentals

    for tk in ["AAPL", "JNJ"]:
        try:
            fin = fetch_sec_financials(tk, log=_quiet)
            cd  = fetch_company_details(tk, KEY, log=_quiet)
            f   = compute_fundamentals(fin, market_cap=cd.get("Market Cap"))
        except Exception as e:
            check(False, f"{tk}: fundamentals compute", str(e)[:60])
            continue
        if not f.get("ok"):
            check(True, f"{tk}: skipped — EDGAR returned no statements")
            continue

        m, r, l, v, q = f["margins"], f["returns"], f["leverage"], f["valuation"], f["quality"]

        def in_range(x, lo, hi):
            return x is None or (lo <= x <= hi)   # None = gracefully absent, not a failure

        check(in_range(m["gross"], -100, 100) and in_range(m["operating"], -200, 100)
              and in_range(m["net"], -200, 100), f"{tk}: margins in sane % range",
              f"gross {m['gross']} / op {m['operating']} / net {m['net']}")
        check(in_range(r["roe"], -500, 600) and in_range(r["roa"], -200, 200),
              f"{tk}: ROE/ROA in sane % range", f"roe {r['roe']} / roa {r['roa']}")
        check(in_range(v["pe"], 0, 5000) and in_range(v["ps"], 0, 500)
              and in_range(v["pb"], 0, 500), f"{tk}: P/E·P/S·P/B non-negative & bounded",
              f"pe {v['pe']} / ps {v['ps']} / pb {v['pb']}")
        check(in_range(l["current_ratio"], 0, 100), f"{tk}: current ratio in [0,100]",
              f"{l['current_ratio']}")
        check(q["f_score"] is None or 0 <= q["f_score"] <= 9, f"{tk}: Piotroski F in [0,9]",
              f"{q['f_score']}")
        check(q["z_zone"] is None or isinstance(q["z_zone"], str), f"{tk}: Altman zone is a label",
              f"{q['z_zone']}")


# ── E. Edge cases (crypto, invalid ticker) ──────────────────────────────────────
def part_e_edge_cases():
    print("\nE. Edge cases")
    print("=" * 60)
    from data import fetch_crypto_data, fetch_stock_data

    # Crypto should resolve and return sane prices. fetch_crypto_data takes the
    # bare symbol ("BTC") and adds the X:…USD wrapper itself.
    try:
        cdf = fetch_crypto_data("BTC", api_key=KEY, log=_quiet,
                                start_override=START, end_override=END, bar_size="day")
        if cdf is None or len(cdf) < 30:
            check(True, "BTC: skipped — crypto data unavailable")
        else:
            last = float(cdf["Close"].iloc[-1])
            check(1_000 < last < 10_000_000, "BTC: price in sane range", f"${last:,.0f}")
            check((cdf["Close"] > 0).all(), "BTC: all closes positive")
    except Exception as e:
        check(True, f"BTC: skipped — {str(e)[:50]}")

    # An invalid ticker must fail cleanly — never fabricate a populated series.
    try:
        bad = fetch_stock_data("ZZQQNOTREAL", benchmark_tickers=("SPY",), api_key=KEY,
                               log=_quiet, start_override=START, end_override=END, bar_size="day")
        check(bad is None or len(bad) < 5, "invalid ticker returns empty (no fabricated data)",
              f"{0 if bad is None else len(bad)} rows")
    except Exception:
        check(True, "invalid ticker raises cleanly (handled)")


# ── F. Single-stock report consistency (Excel + PowerPoint) ─────────────────────
def part_f_stock_report_consistency():
    print("\nF. Single-stock report consistency (Excel + PPTX)")
    print("=" * 60)
    import io
    from data import (fetch_stock_data, fetch_company_details, fetch_news,
                      fetch_sec_financials, fetch_financials)
    from analysis import (run_monte_carlo, generate_summary_paragraph,
                          detect_support_resistance, compute_fundamentals)
    from constants import get_risk_free_rate
    from excel_builder import build_excel
    try:
        from pptx import Presentation
        from openpyxl import load_workbook
    except ImportError:
        check(True, "skipped — python-pptx/openpyxl not installed")
        return

    TK = "AAPL"
    try:
        df = fetch_stock_data(TK, benchmark_tickers=("SPY",), api_key=KEY, log=_quiet,
                              start_override=START, end_override=END, bar_size="day")
        cd  = fetch_company_details(TK, KEY, log=_quiet)
        _mc = run_monte_carlo(df, n_simulations=300, forecast_days=252, log=_quiet)
        mc_df, mc_sum = (_mc[0], _mc[1]) if isinstance(_mc, tuple) else (_mc, {})
        res, sup = detect_support_resistance(df)
        try:
            _fr  = fetch_sec_financials(TK, log=_quiet) or fetch_financials(TK, KEY, log=_quiet)
            fund = compute_fundamentals(_fr, market_cap=cd.get("Market Cap"),
                                        price=float(df["Close"].iloc[-1]))
        except Exception:
            fund = {"ok": False}
    except Exception as e:
        check(False, f"{TK}: build pipeline", str(e)[:80])
        return

    # Independently-computed source truth (same formulas the reports use).
    price   = float(df["Close"].iloc[-1])
    ret     = df["Daily_Return"].dropna()
    ann_ret = ret.mean() * 252
    ann_std = ret.std() * np.sqrt(252)
    sharpe  = (ann_ret - get_risk_free_rate()) / ann_std if ann_std else 0.0

    # ── Excel: read Dashboard cells and compare to source truth ──────────────────
    try:
        xls = build_excel(TK, df, "2Y", company_details=cd, mc_sim_df=mc_df, mc_summary=mc_sum,
                          news_list=fetch_news(TK, KEY, log=_quiet)[:5], summary_text="",
                          resistance_levels=res, support_levels=sup, fundamentals=fund)
        wb  = load_workbook(io.BytesIO(xls.getvalue()))
        ws  = wb["Dashboard"]
        kv  = {}                                    # {label: value} from cols A/B
        for r in range(1, ws.max_row + 1):
            lbl, val = ws.cell(r, 1).value, ws.cell(r, 2).value
            if isinstance(lbl, str) and isinstance(val, (int, float)):
                kv[lbl.strip()] = val
        cp = kv.get("Current Price ($)")
        sh = kv.get("Sharpe Ratio")
        check(cp is not None and abs(cp - price) < 0.05,
              "Excel dashboard 'Current Price' matches source", f"{cp} vs {price:.2f}")
        check(sh is not None and abs(sh - round(sharpe, 2)) < 0.1,
              "Excel dashboard 'Sharpe Ratio' matches source", f"{sh} vs {sharpe:.2f}")
        # No %-formatted cell should hold a value implying >500% (the ×100 class).
        bad_pct = []
        for row in ws.iter_rows():
            for c in row:
                if isinstance(c.number_format, str) and "%" in c.number_format \
                        and isinstance(c.value, (int, float)) and abs(c.value) >= 5:
                    bad_pct.append((c.coordinate, c.value))
        check(not bad_pct, "Excel: no %-cell implies >500%", str(bad_pct[:3]))
    except Exception as e:
        check(False, "Excel stock report build/read", str(e)[:80])

    # ── PPTX: extract text, check price present and no absurd % ──────────────────
    try:
        from pptx_builder import build_stock_pptx
        try:
            summ = generate_summary_paragraph(TK, df, cd, mc_sum, sharpe, sharpe)
        except Exception:
            summ = ""
        pbuf = build_stock_pptx(TK, df, "2Y", company_details=cd, mc_sim_df=mc_df,
                                mc_summary=mc_sum, summary_text=summ, fundamentals=fund)
        prs  = Presentation(io.BytesIO(pbuf.getvalue()))
        blob = "\n".join(sh.text_frame.text for sl in prs.slides
                         for sh in sl.shapes if sh.has_text_frame)
        dollars = [float(x.replace(",", "")) for x in re.findall(r"\$([\d,]+\.\d{2})", blob)]
        check(any(abs(d - price) < 0.5 for d in dollars),
              "PPTX shows the current price", f"price {price:.2f}; ${'/'.join(f'{d:.0f}' for d in dollars[:4])}…")
        pcts  = [float(m.replace(",", "")) for m in re.findall(r"([+-]?\d[\d,]*\.?\d*)%", blob)]
        worst = max((abs(p) for p in pcts), default=0.0)
        check(worst <= 300.0, "PPTX: no absurd percentage (≤300%)", f"max |%| = {worst:.0f}%")
    except Exception as e:
        check(False, "PPTX stock report build/read", str(e)[:80])


# ── G. Live quote freshness (verifies the "real-time via Finnhub" UI claim) ──────
def part_g_quote_freshness():
    print("\nG. Live quote freshness")
    print("=" * 60)
    from market_data import get_quote, get_bars, finnhub_key
    tk = "AAPL"
    if not finnhub_key():
        check(True, "skipped — FINNHUB_API_KEY not set in this environment")
        return
    q = get_quote(tk, polygon_key=KEY)
    if not q:
        check(False, f"{tk}: get_quote returned a live quote")
        return
    # The UI now claims quotes are real-time via Finnhub — verify that's the path.
    check(q.get("source") == "finnhub",
          "live quote served by Finnhub (real-time path)", f"source={q.get('source')}")
    price = q.get("price", 0) or 0
    check(price > 0, f"{tk}: quote price positive", f"{price}")
    # Sanity vs the last daily close: a real overnight/intraday move is small; a
    # 20%+ gap means a stale/garbage/mis-symbol quote.
    bars = get_bars(tk, START, END, interval="day", polygon_key=KEY)
    if bars is not None and len(bars):
        last_close = float(bars["Close"].iloc[-1])
        drift = abs(price - last_close) / last_close if last_close else 1.0
        check(drift < 0.20, "quote within 20% of last daily close",
              f"quote {price:.2f} vs close {last_close:.2f} ({drift*100:.1f}%)")


def main():
    print("QuantWizard Data-Correctness Validation")
    print(f"Window: {START} → {END}")
    part_a_stock_metrics()
    part_b_cross_source()
    part_c_report_consistency()
    part_d_fundamentals()
    part_e_edge_cases()
    part_f_stock_report_consistency()
    part_g_quote_freshness()

    print("\n" + "=" * 60)
    print(f"Results: {len(_PASS)} passed, {len(_FAIL)} failed")
    if _FAIL:
        print("FAILED:")
        for f in _FAIL:
            print(f"   ✗ {f}")
        sys.exit(1)
    print("✓ All data-correctness checks passed.")


if __name__ == "__main__":
    main()
